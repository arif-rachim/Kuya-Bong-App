#!/usr/bin/env python3
"""Build the 3-slide client-facing prototype-update deck (+ cover).

Follows docs/Kuya_Bong_Slides-Guidance.pdf:
  Cover (not counted) + Slide 1 business situation (4 figures) +
  Slide 2 one app / two roles + Slide 3 simple investment options.

Run: python3 scripts/build_summary_deck.py
Output: docs/Kuya_Bong_Prototype_Update_3Slide.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SHOTS = os.path.join(ROOT, "docs/assets/shots")
GEN = os.path.join(ROOT, "docs/assets/gen")
OUT = os.path.join(ROOT, "docs/Kuya_Bong_Prototype_Update_3Slide_v2.pptx")

GREEN = RGBColor(0x1E, 0x9E, 0x3A)
GREEN_DK = RGBColor(0x12, 0x5E, 0x26)
GREEN_SOFT = RGBColor(0xE8, 0xF5, 0xEC)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
INK = RGBColor(0x1B, 0x2A, 0x22)
SLATE = RGBColor(0x53, 0x63, 0x5B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xF4, 0xF8, 0xF5)
AMBER_BG = RGBColor(0xFB, 0xF3, 0xDD)
AMBER_TX = RGBColor(0x8A, 0x5E, 0x0A)
FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def rect(s, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE, line_w=1.0):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def soft_shadow(sp, alpha=22000):
    spPr = sp._element.spPr
    el = spPr.makeelement(qn('a:effectLst'), {})
    sh = el.makeelement(qn('a:outerShdw'),
                        {'blurRad': '90000', 'dist': '40000', 'dir': '5400000', 'rotWithShape': '0'})
    clr = sh.makeelement(qn('a:srgbClr'), {'val': '1B2A22'})
    a = clr.makeelement(qn('a:alpha'), {'val': str(alpha)}); clr.append(a)
    sh.append(clr); el.append(sh); spPr.append(el)


def set_alpha(sp, alpha):
    sf = sp.fill.fore_color._xFill.find(qn('a:srgbClr'))
    if sf is not None:
        a = sf.makeelement(qn('a:alpha'), {'val': str(alpha)}); sf.append(a)


def P(text, size, bold=False, color=INK, font=FONT):
    return [(text, size, bold, color, font)]


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.05):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, m, 0)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, b, col, fnt) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = col; r.font.name = fnt
    return tb


def pic_cover(s, path, x, y, w, h):
    if not os.path.exists(path):
        rect(s, x, y, w, h, CARD); return None
    iw, ih = Image.open(path).size
    box_ar = w / h; img_ar = iw / ih
    pic = s.shapes.add_picture(path, x, y, w, h)
    if img_ar > box_ar:
        c = (1 - box_ar / img_ar) / 2; pic.crop_left = c; pic.crop_right = c
    else:
        c = (1 - img_ar / box_ar) / 2; pic.crop_top = c; pic.crop_bottom = c
    return pic


def kicker(s, x, y, text, color=GREEN):
    rect(s, x, y + Emu(22000), Emu(300000), Emu(66000), color)
    txt(s, x + Emu(380000), y, Inches(9), Inches(0.4), [P(text.upper(), 13, True, color)])


def chip(s, x, y, w, head, body):
    h = Inches(1.02)
    c = rect(s, x, y, w, h, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    soft_shadow(c, 16000)
    rect(s, x + Inches(0.16), y + Inches(0.2), Inches(0.12), h - Inches(0.4), GREEN,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x + Inches(0.42), y + Inches(0.13), w - Inches(0.6), h - Inches(0.2),
        [P(head, 12.5, True, INK), P(body, 9.5, False, SLATE)], space_after=2, line_spacing=1.0,
        anchor=MSO_ANCHOR.MIDDLE)


# ===========================================================================
# COVER (not one of the 3 slides)
# ===========================================================================
def cover():
    s = slide()
    bg(s, WHITE)
    pic_cover(s, os.path.join(GEN, "hero-cover.png"), Inches(7.1), 0, Inches(6.233), SH)
    rect(s, Inches(7.05), 0, Inches(0.08), SH, GREEN)
    rect(s, Inches(0.7), Inches(0.72), Inches(0.5), Inches(0.5), GREEN, shape=MSO_SHAPE.OVAL)
    txt(s, Inches(1.4), Inches(0.74), Inches(5), Inches(0.5), [P("REALIEF EXPERT", 14, True, GREEN_DK)])
    txt(s, Inches(0.7), Inches(2.15), Inches(6.1), Inches(1.6),
        [P("Kuya Bong App", 46, True, INK),
         P("Prototype Update & Feedback Session", 19, True, GREEN_DK)],
        line_spacing=1.0, space_after=8)
    txt(s, Inches(0.72), Inches(3.85), Inches(5.9), Inches(1.0),
        [P("A simple mobile app to help manage bookings, services, packages, and product follow-ups.",
           14, False, SLATE)], line_spacing=1.2)
    rect(s, Inches(0.7), Inches(5.2), Inches(2.6), Inches(0.02), GREEN)
    txt(s, Inches(0.7), Inches(5.4), Inches(6), Inches(1.3),
        [P("Prepared by:  Analyst & Sales Team", 13, True, INK),
         P("Meeting date:  19 June 2026", 13, False, SLATE),
         P("Expert Care, Real Relief  ·  Physiotherapy & Chiropractic", 11, False, SLATE)],
        space_after=5)


# ===========================================================================
# SLIDE 1 — Current business situation (4 figures)
# ===========================================================================
def slide1():
    s = slide()
    bg(s, WHITE)
    kicker(s, Inches(0.6), Inches(0.45), "Slide 1 · Today")
    txt(s, Inches(0.56), Inches(0.82), Inches(12), Inches(0.8),
        [P("Understanding Kuya Bong's current business situation", 27, True, INK)])
    figs = [
        ("fig-clinics.png", "Two clinics, one therapist",
         "Kuya Bong moves between Clinic A & Clinic B to serve booked patients — schedules must avoid time conflicts."),
        ("fig-services-products.png", "Two kinds of business",
         "Appointment-based services (Physiotherapy & Massage, Grounding Therapy) and herbal / supplement product sales."),
        ("fig-booking.png", "Booking & package pain",
         "Patients can't see free slots, may pick the wrong clinic or time, and can't track paid package sessions."),
        ("fig-followup.png", "Product follow-up",
         "Hard to remember what product was sold, when, and when to follow up before it runs out."),
    ]
    x0, y0 = Inches(0.6), Inches(1.85)
    cw, ch = Inches(6.0), Inches(2.55)
    gx, gy = Inches(0.18), Inches(0.18)
    for i, (img, h, msg) in enumerate(figs):
        col, row = i % 2, i // 2
        x = x0 + col * (cw + gx); y = y0 + row * (ch + gy)
        c = rect(s, x, y, cw, ch, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        soft_shadow(c, 14000)
        # image thumb
        th = Inches(2.15)
        tw = th
        pic_cover(s, os.path.join(GEN, img), x + Inches(0.18), y + Inches(0.2), tw, th)
        # text
        txt(s, x + Inches(2.55), y + Inches(0.28), cw - Inches(2.75), ch - Inches(0.5),
            [P(f"Figure {i+1}", 10, True, GREEN),
             P(h, 16, True, INK),
             [(msg, 11.5, False, SLATE, FONT)]],
            space_after=5, line_spacing=1.08, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, "Current situation")


# ===========================================================================
# SLIDE 2 — One app, two roles
# ===========================================================================
def slide2():
    s = slide()
    bg(s, WHITE)
    kicker(s, Inches(0.6), Inches(0.45), "Slide 2 · The solution")
    txt(s, Inches(0.56), Inches(0.82), Inches(12), Inches(0.8),
        [P("One app serving two main roles", 27, True, INK)])

    # center phone
    cx = SW // 2
    ph_h = Inches(3.85)
    ar = 412 / 892
    ph_w = Emu(int(ph_h * ar))
    px = Emu(int(cx - ph_w / 2))
    top = Inches(1.95)
    fr = rect(s, px - Emu(55000), top - Emu(55000), ph_w + Emu(110000), ph_h + Emu(110000),
              WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    soft_shadow(fr, 26000)
    wp = os.path.join(SHOTS, "welcome.png")
    if os.path.exists(wp):
        s.shapes.add_picture(wp, px, top, ph_w, ph_h)
    else:
        rect(s, px, top, ph_w, ph_h, CARD)
    txt(s, Emu(int(cx)) - Inches(1.8), top + ph_h + Inches(0.16), Inches(3.6), Inches(0.7),
        [P("Kuya Bong App", 15, True, GREEN_DK),
         P("One mobile app · role-based access", 10.5, False, SLATE)],
        align=PP_ALIGN.CENTER, space_after=1)

    # role headers
    txt(s, Inches(0.6), Inches(1.55), Inches(4.4), Inches(0.4),
        [[("●  ", 12, True, GREEN, FONT), ("KUYA BONG / ADMIN", 13, True, GREEN_DK, FONT)]])
    txt(s, Inches(8.33), Inches(1.55), Inches(4.4), Inches(0.4),
        [[("●  ", 12, True, TEAL, FONT), ("PATIENT", 13, True, TEAL, FONT)]], align=PP_ALIGN.RIGHT)

    admin = [
        "Manage clinics, services & therapists",
        "Publish availability (clinic · service · time)",
        "View, approve, move & create bookings",
        "Packages: assign, track use, balance & expiry",
        "Products, price & product follow-up list",
        "Manage cancellation reasons",
    ]
    patient = [
        "Register, log in & reset password",
        "Book: service → clinic → date → time",
        "Reschedule or cancel (with a reason)",
        "View appointments & booking confirmation",
        "View package balance & expiry",
        "Family link: spouse & children",
    ]
    yy = Inches(2.0)
    step = Inches(0.66)
    lw = Inches(3.95)
    for t in admin:
        _feat(s, Inches(0.6), yy, lw, t, GREEN, PP_ALIGN.LEFT)
        yy += step
    yy = Inches(2.0)
    for t in patient:
        _feat(s, Inches(8.78), yy, lw, t, TEAL, PP_ALIGN.LEFT)
        yy += step

    # key message
    band = rect(s, Inches(0.6), Inches(6.5), Inches(12.13), Inches(0.62), GREEN_SOFT,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(0.9), Inches(6.55), Inches(11.6), Inches(0.5),
        [[("Key message:  ", 13, True, GREEN_DK, FONT),
          ("We are not building two apps — it is one app with different access depending on the user role.",
           13, False, GREEN_DK, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, "One app, two roles")


def _feat(s, x, y, w, text, color, align):
    h = Inches(0.56)
    c = rect(s, x, y, w, h, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x + Inches(0.12), y + Inches(0.14), Inches(0.1), h - Inches(0.28), color,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x + Inches(0.34), y, w - Inches(0.5), h, [P(text, 11, True, INK)],
        anchor=MSO_ANCHOR.MIDDLE)


# ===========================================================================
# SLIDE 3 — Simple investment options
# ===========================================================================
def slide3():
    s = slide()
    bg(s, WHITE)
    kicker(s, Inches(0.6), Inches(0.45), "Slide 3 · Investment")
    txt(s, Inches(0.56), Inches(0.82), Inches(12), Inches(0.8),
        [P("Simple investment options to start using the app", 27, True, INK)])

    # estimate banner
    band = rect(s, Inches(0.6), Inches(1.6), Inches(12.13), Inches(0.55), AMBER_BG,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(0.85), Inches(1.64), Inches(11.7), Inches(0.5),
        [[("All amounts are estimates for discussion.  ", 11.5, True, AMBER_TX, FONT),
          ("Store, hosting, domain, SMS & email are external costs that vary by provider, country & usage.",
           11.5, False, AMBER_TX, FONT)]], anchor=MSO_ANCHOR.MIDDLE)

    kickoff = {
        "tag": "STARTING COST", "name": "Kick-off Cost", "color": GREEN,
        "sub": "One-time first-year setup — free / starter tiers where possible.",
        "items": [
            ("Apple Store registration", "USD 100"),
            ("Play Store / Google registration", "USD 100"),
            ("Database (free / starter tier)", "USD 0"),
            ("Domain name", "USD 15"),
        ],
        "total_label": "Total Starting Cost", "total_val": "USD 215  (AED 790)",
        "note": None,
    }
    maint = {
        "tag": "ANNUALLY", "name": "Maintenance Cost", "color": TEAL,
        "sub": "Recurring yearly cost from next year onward.",
        "items": [
            ("Domain name renewal", "USD 15"),
            ("Apple developer account (yearly)", "USD 100"),
            ("Google Play (one-time — no yearly fee)", "USD 0"),
            ("Database (within free / starter tier)", "USD 0"),
        ],
        "total_label": "Total Annual Maintenance", "total_val": "USD 115  (AED 422)",
        "note": "Database may no longer be free if capacity grows significantly — considered unlikely at this stage.",
    }
    cw = Inches(6.0); ch = Inches(4.6)
    _cost_card(s, Inches(0.6), Inches(2.35), cw, ch, kickoff)
    _cost_card(s, Inches(6.73), Inches(2.35), cw, ch, maint)
    txt(s, Inches(0.6), Inches(7.08), Inches(12), Inches(0.4),
        [P("* Estimates — to be confirmed. USD→AED at 1 USD ≈ 3.67 AED. Store, hosting & domain are external costs that vary by provider, country & account type.",
           9, False, SLATE)])


def _cost_card(s, x, y, w, h, o):
    c = rect(s, x, y, w, h, WHITE, line=RGBColor(0xE0, 0xEA, 0xE3), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    soft_shadow(c, 16000)
    # header bar
    hb = rect(s, x, y, w, Inches(0.95), o["color"], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x + Inches(0.35), y + Inches(0.14), w - Inches(0.7), Inches(0.7),
        [P(o["tag"], 10, True, WHITE),
         [(o["name"], 21, True, WHITE, FONT)]], space_after=1)
    txt(s, x + Inches(0.35), y + Inches(1.02), w - Inches(0.7), Inches(0.45),
        [P(o["sub"], 11, False, SLATE)], line_spacing=1.0)
    # items: label on the left, amount right-aligned
    yy = y + Inches(1.55)
    for it, amount in o["items"]:
        rect(s, x + Inches(0.38), yy + Inches(0.06), Inches(0.13), Inches(0.13), o["color"],
             shape=MSO_SHAPE.OVAL)
        txt(s, x + Inches(0.66), yy - Inches(0.02), w - Inches(2.1), Inches(0.4),
            [P(it, 11.5, True, INK)], line_spacing=1.0, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x + w - Inches(1.55), yy - Inches(0.02), Inches(1.25), Inches(0.4),
            [P(amount, 11.5, True, SLATE)], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        yy += Inches(0.46)
    # total band (same position on both cards)
    pb_y = y + h - Inches(1.15)
    rect(s, x + Inches(0.3), pb_y, w - Inches(0.6), Inches(0.62), GREEN_SOFT,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x + Inches(0.5), pb_y, Inches(2.7), Inches(0.62),
        [P(o["total_label"], 12, True, INK)], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    txt(s, x + w - Inches(3.4), pb_y, Inches(2.9), Inches(0.62),
        [P(o["total_val"], 15, True, GREEN_DK)], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    if o["note"]:
        txt(s, x + Inches(0.4), pb_y + Inches(0.7), w - Inches(0.8), Inches(0.5),
            [[("Note:  ", 9.5, True, AMBER_TX, FONT), (o["note"], 9.5, False, SLATE, FONT)]],
            line_spacing=1.02)


def footer(s, label):
    txt(s, Inches(0.6), Inches(7.12), Inches(8), Inches(0.3),
        [P("Kuya Bong App  ·  Prototype Update & Feedback Session", 9, False, SLATE)])
    txt(s, Inches(10.5), Inches(7.12), Inches(2.23), Inches(0.3),
        [P(label, 9, False, SLATE)], align=PP_ALIGN.RIGHT)


cover()
slide1()
slide2()
slide3()
prs.save(OUT)
print("Saved", OUT, "with", len(prs.slides._sldIdLst), "slides (cover + 3)")
