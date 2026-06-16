#!/usr/bin/env python3
"""Build the Realief Expert / Kuya Bong product & cost presentation (PPTX).

Run: python3 scripts/build_deck.py
Output: docs/Realief_Expert_Kuya_Bong_Presentation.pptx
Assets: docs/assets/shots (live-app screenshots), docs/assets/gen (AI images).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SHOTS = os.path.join(ROOT, "docs/assets/shots")
GEN = os.path.join(ROOT, "docs/assets/gen")
OUT = os.path.join(ROOT, "docs/Realief_Expert_Kuya_Bong_Presentation.pptx")

# ---- palette -------------------------------------------------------------
GREEN = RGBColor(0x1E, 0x9E, 0x3A)      # brand
GREEN_DK = RGBColor(0x12, 0x5E, 0x26)   # deep green
GREEN_SOFT = RGBColor(0xE8, 0xF5, 0xEC) # light wash
INK = RGBColor(0x1B, 0x2A, 0x22)        # near-black text
SLATE = RGBColor(0x53, 0x63, 0x5B)      # muted text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xF4, 0xF8, 0xF5)
CLINIC_A = RGBColor(0x1E, 0x9E, 0x3A)
CLINIC_B = RGBColor(0x0E, 0x7C, 0x7B)
AMBER = RGBColor(0xC8, 0x86, 0x12)

FONT = "Calibri"
FONT_H = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---- low-level helpers ---------------------------------------------------
def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def rect(s, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.05):
    """runs: list of paragraphs; each paragraph is list of (text, size, bold, color, font)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, b, col, fnt) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(sz)
            r.font.bold = b
            r.font.color.rgb = col
            r.font.name = fnt
    return tb


def P(text, size, bold=False, color=INK, font=FONT):
    return [(text, size, bold, color, font)]


def pic_cover(s, path, x, y, w, h):
    """Place an image cropping to fill the (w,h) box (object-fit: cover)."""
    if not os.path.exists(path):
        rect(s, x, y, w, h, CARD)
        return None
    iw, ih = Image.open(path).size
    box_ar = w / h
    img_ar = iw / ih
    pic = s.shapes.add_picture(path, x, y, w, h)
    if img_ar > box_ar:
        crop = (1 - box_ar / img_ar) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    else:
        crop = (1 - img_ar / box_ar) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop
    return pic


def phone(s, name, cx, top, ph_h):
    """Place a phone screenshot centered at cx, height ph_h, on a white card with caption-free frame."""
    path = os.path.join(SHOTS, name)
    # device aspect 412x892
    ar = 412 / 892
    ph_w = Emu(int(ph_h * ar))
    x = Emu(int(cx - ph_w / 2))
    # frame
    fr = rect(s, x - Emu(60000), top - Emu(60000), ph_w + Emu(120000), ph_h + Emu(120000),
              WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _soft_shadow(fr)
    if os.path.exists(path):
        s.shapes.add_picture(path, x, top, ph_w, ph_h)
    else:
        rect(s, x, top, ph_w, ph_h, CARD)
    return ph_w


def _soft_shadow(sp):
    spPr = sp._element.spPr
    el = spPr.makeelement(qn('a:effectLst'), {})
    sh = el.makeelement(qn('a:outerShdw'),
                        {'blurRad': '90000', 'dist': '40000', 'dir': '5400000', 'rotWithShape': '0'})
    clr = sh.makeelement(qn('a:srgbClr'), {'val': '1B2A22'})
    alpha = clr.makeelement(qn('a:alpha'), {'val': '22000'})
    clr.append(alpha)
    sh.append(clr)
    el.append(sh)
    spPr.append(el)


def kicker(s, x, y, text, color=GREEN):
    rect(s, x, y + Emu(20000), Emu(300000), Emu(70000), color)
    txt(s, x + Emu(380000), y, Inches(6), Inches(0.4),
        [P(text.upper(), 13, True, color)])


def footer(s, n, label):
    txt(s, Inches(0.55), Inches(7.05), Inches(8), Inches(0.35),
        [P("Realief Expert  ·  Expert Care, Real Relief", 9, False, SLATE)])
    txt(s, Inches(10.5), Inches(7.05), Inches(2.3), Inches(0.35),
        [P(f"{label}  ·  {n:02d}", 9, False, SLATE)], align=PP_ALIGN.RIGHT)


# ===========================================================================
# SLIDE BUILDERS
# ===========================================================================
PAGE = {"n": 0}


def n():
    PAGE["n"] += 1
    return PAGE["n"]


def title_slide():
    s = slide()
    bg(s, WHITE)
    # left band
    rect(s, 0, 0, Inches(7.6), SH, WHITE)
    pic_cover(s, os.path.join(GEN, "hero-cover.png"), Inches(7.0), 0, Inches(6.333), SH)
    # green seam
    rect(s, Inches(6.95), 0, Inches(0.08), SH, GREEN)
    rect(s, Inches(0.7), Inches(0.7), Inches(0.55), Inches(0.55), GREEN, shape=MSO_SHAPE.OVAL)
    txt(s, Inches(1.45), Inches(0.72), Inches(5), Inches(0.6),
        [P("REALIEF EXPERT", 16, True, GREEN_DK)])
    txt(s, Inches(0.7), Inches(2.3), Inches(6.0), Inches(2.2),
        [P("Kuya Bong", 30, True, SLATE),
         P("Clinic Booking &", 44, True, INK),
         P("Management App", 44, True, GREEN_DK)], line_spacing=1.0, space_after=2)
    txt(s, Inches(0.72), Inches(4.65), Inches(5.9), Inches(1.2),
        [P("Product walkthrough & operating-cost overview for a single iOS + Android app serving Patients and the Administrator across two clinics.",
           15, False, SLATE)], line_spacing=1.15)
    rect(s, Inches(0.7), Inches(6.05), Inches(3.0), Inches(0.02), GREEN)
    txt(s, Inches(0.7), Inches(6.2), Inches(6), Inches(0.8),
        [P("Expert Care, Real Relief", 14, True, GREEN),
         P("Physiotherapy & Chiropractic  ·  Blueprint v0.3  ·  2026", 11, False, SLATE)],
        space_after=2)


def agenda_slide():
    s = slide()
    bg(s, WHITE)
    kicker(s, Inches(0.7), Inches(0.6), "Agenda")
    txt(s, Inches(0.66), Inches(1.0), Inches(11), Inches(0.9),
        [P("What's inside", 34, True, INK)])
    items = [
        ("01", "The problem", "Manual scheduling, package & follow-up tracking pain points."),
        ("02", "The solution", "One app, two roles, two clinics — built around the blueprint."),
        ("03", "Patient experience", "Onboarding, booking, appointments, packages, family."),
        ("04", "Admin experience", "Calendar, bookings, packages, products, master data."),
        ("05", "Requirements coverage", "How screens map to the confirmed user stories & rules."),
        ("06", "Operating costs", "Stores, Apple Pay, servers, database, domain — year-1 view."),
    ]
    x0, y0 = Inches(0.7), Inches(2.0)
    cw, ch = Inches(5.9), Inches(1.45)
    gx, gy = Inches(0.25), Inches(0.2)
    for i, (num, h, d) in enumerate(items):
        col, row = i % 2, i // 2
        x = x0 + col * (cw + gx)
        y = y0 + row * (ch + gy)
        c = rect(s, x, y, cw, ch, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, y, Inches(0.12), ch, GREEN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, x + Inches(0.35), y + Inches(0.18), Inches(1.1), Inches(1),
            [P(num, 30, True, GREEN_SOFT_DARK := GREEN)])
        txt(s, x + Inches(1.5), y + Inches(0.22), cw - Inches(1.7), Inches(1.1),
            [P(h, 18, True, INK), P(d, 12, False, SLATE)], space_after=3, line_spacing=1.05)
    footer(s, n(), "Agenda")


def section_slide(title, sub, img, kick):
    s = slide()
    bg(s, GREEN_DK)
    pic_cover(s, os.path.join(GEN, img), Inches(7.2), 0, Inches(6.133), SH)
    # dark overlay gradient via translucent rect
    ov = rect(s, Inches(7.2), 0, Inches(6.133), SH, GREEN_DK)
    _set_alpha(ov, 38000)
    rect(s, 0, 0, Inches(7.2), SH, GREEN_DK)
    rect(s, Inches(0.7), Inches(2.7), Inches(1.4), Inches(0.08), GREEN)
    txt(s, Inches(0.7), Inches(2.05), Inches(6), Inches(0.5),
        [P(kick.upper(), 14, True, RGBColor(0x9F, 0xE6, 0xB0))])
    txt(s, Inches(0.66), Inches(2.95), Inches(6.2), Inches(2),
        [P(title, 40, True, WHITE)], line_spacing=1.0)
    txt(s, Inches(0.7), Inches(4.6), Inches(5.9), Inches(1.6),
        [P(sub, 15, False, RGBColor(0xD7, 0xEC, 0xDD))], line_spacing=1.2)


def _set_alpha(sp, alpha):
    sf = sp.fill.fore_color._xFill.find(qn('a:srgbClr'))
    if sf is not None:
        a = sf.makeelement(qn('a:alpha'), {'val': str(alpha)})
        sf.append(a)


def bullets_slide(title, kick, bullets, note=None):
    s = slide()
    bg(s, WHITE)
    kicker(s, Inches(0.7), Inches(0.55), kick)
    txt(s, Inches(0.66), Inches(0.95), Inches(11.5), Inches(0.9), [P(title, 30, True, INK)])
    nb = len(bullets)
    y = Inches(1.95)
    step = Inches(0.84) if nb <= 5 else Inches(0.72)
    for head, body in bullets:
        rect(s, Inches(0.75), y + Inches(0.07), Inches(0.16), Inches(0.16), GREEN, shape=MSO_SHAPE.OVAL)
        txt(s, Inches(1.15), y - Inches(0.02), Inches(11.2), Inches(0.9),
            [[(head + "  ", 16, True, INK, FONT), (body, 14, False, SLATE, FONT)]], line_spacing=1.1)
        y += step
    if note:
        rect(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.55), GREEN_SOFT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, Inches(0.95), Inches(6.43), Inches(11.5), Inches(0.45),
            [P(note, 12, True, GREEN_DK)], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, n(), kick)


def feature_slide(kick, title, shots, points):
    """shots: list of (filename, caption). points: list of (head, body)."""
    s = slide()
    bg(s, WHITE)
    rect(s, 0, 0, Inches(4.55), SH, GREEN_SOFT)
    kicker(s, Inches(0.55), Inches(0.55), kick)
    txt(s, Inches(0.52), Inches(0.95), Inches(3.9), Inches(1.4), [P(title, 26, True, INK)], line_spacing=0.98)
    y = Inches(2.45)
    for head, body in points:
        rect(s, Inches(0.6), y + Inches(0.05), Inches(0.14), Inches(0.14), GREEN, shape=MSO_SHAPE.OVAL)
        txt(s, Inches(0.92), y - Inches(0.03), Inches(3.45), Inches(1.2),
            [P(head, 14, True, GREEN_DK), P(body, 11.5, False, SLATE)], space_after=2, line_spacing=1.05)
        y += Inches(0.96)
    # phones on the right
    area_x = Inches(4.55)
    area_w = SW - area_x
    ph_h = Inches(5.1)
    top = Inches(1.15)
    nshots = len(shots)
    if nshots == 1:
        centers = [area_x + area_w // 2]
    elif nshots == 2:
        centers = [area_x + area_w * 0.3, area_x + area_w * 0.72]
    else:
        centers = [area_x + area_w * (i + 0.5) / nshots for i in range(nshots)]
        ph_h = Inches(4.7)
        top = Inches(1.25)
    for (fn, cap), cx in zip(shots, centers):
        pw = phone(s, fn, int(cx), top, ph_h)
        txt(s, Emu(int(cx) - Inches(1.4)), top + ph_h + Inches(0.18), Inches(2.8), Inches(0.5),
            [P(cap, 11, True, SLATE)], align=PP_ALIGN.CENTER)
    footer(s, n(), kick)


def table_slide(title, kick, headers, rows, note=None, col_w=None, highlight_last=False):
    s = slide()
    bg(s, WHITE)
    kicker(s, Inches(0.7), Inches(0.55), kick)
    txt(s, Inches(0.66), Inches(0.95), Inches(11.5), Inches(0.9), [P(title, 28, True, INK)])
    ncol = len(headers)
    nrow = len(rows) + 1
    tx, ty = Inches(0.7), Inches(1.95)
    tw = Inches(11.95)
    th = Inches(4.2) if not note else Inches(3.9)
    gtbl = s.shapes.add_table(nrow, ncol, tx, ty, tw, th).table
    if col_w:
        total = sum(col_w)
        for i, w in enumerate(col_w):
            gtbl.columns[i].width = Emu(int(int(tw) * w / total))
    # header
    for j, hd in enumerate(headers):
        c = gtbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = GREEN_DK
        _cell_text(c, hd, 12.5, True, WHITE, PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    for i, row in enumerate(rows):
        last = highlight_last and i == len(rows) - 1
        for j, val in enumerate(row):
            c = gtbl.cell(i + 1, j)
            c.fill.solid()
            c.fill.fore_color.rgb = (GREEN_SOFT if last else (CARD if i % 2 == 0 else WHITE))
            bold = last or j == 0
            col = GREEN_DK if last else (INK if j == 0 else SLATE)
            _cell_text(c, val, 11.5, bold, col, PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
    if note:
        rect(s, Inches(0.7), Inches(6.35), Inches(11.95), Inches(0.55), GREEN_SOFT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, Inches(0.95), Inches(6.43), Inches(11.5), Inches(0.45),
            [P(note, 11.5, True, GREEN_DK)], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, n(), kick)


def _cell_text(cell, text, size, bold, color, align):
    cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = FONT


def closing_slide():
    s = slide()
    bg(s, GREEN_DK)
    pic_cover(s, os.path.join(GEN, "hero-cover.png"), 0, 0, SW, SH)
    ov = rect(s, 0, 0, SW, SH, GREEN_DK)
    _set_alpha(ov, 60000)
    txt(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(2),
        [P("Expert Care, Real Relief", 18, True, RGBColor(0x9F, 0xE6, 0xB0)),
         P("One app. Two roles. Two clinics.", 40, True, WHITE)], space_after=8)
    txt(s, Inches(0.92), Inches(4.4), Inches(11), Inches(1.5),
        [P("A working prototype today — booking, packages, family sharing, products and "
           "master data — ready for backend integration and store release.", 15, False, RGBColor(0xE2, 0xF1, 0xE7))],
        line_spacing=1.2)
    rect(s, Inches(0.95), Inches(5.7), Inches(2.6), Inches(0.03), GREEN)
    txt(s, Inches(0.95), Inches(5.85), Inches(11), Inches(0.5),
        [P("Live demo: realief-expert.netlify.app   ·   github.com/arif-rachim/Kuya-Bong-App", 12, True, WHITE)])


# ===========================================================================
# CONTENT
# ===========================================================================
title_slide()
agenda_slide()

# --- problem / solution ---
section_slide(
    "The problem today",
    "Kuya Bong is a well-known physiotherapist & chiropractor running two clinics. "
    "Scheduling, package tracking and product follow-ups are all handled manually.",
    "concept-clinic.png", "Context")
bullets_slide(
    "Why a dedicated app", "The problem",
    [("Booking is manual & error-prone.", "Patients don't know open slots or which clinic; phone/DM requests are easy to lose."),
     ("Packages are hard to track.", "Sessions used, balance left, expiry and which family member used a session."),
     ("Clinic confusion.", "Patients arrive at the wrong clinic (A vs B) or forget appointment details."),
     ("Product follow-up is reactive.", "No record of last herbal/supplement purchase to prompt timely re-orders."),
     ("No single source of truth.", "App, phone and walk-in bookings live in different places.")],
    note="Source: Product Concept & Solution Blueprint v0.3 — Sections 3, 6 & user stories.")
section_slide(
    "The solution",
    "A single downloadable app — published to the App Store and Google Play — serving both "
    "Patients and the Administrator through role-based access.",
    "concept-booking.png", "Proposed solution")
bullets_slide(
    "One app, two roles, two clinics", "The solution",
    [("Single mobile app.", "No separate patient/admin apps — role decides which screens appear."),
     ("Service-driven booking.", "Pick a service first; its duration drives valid start times and end time."),
     ("Conflict-free scheduling.", "Blocks therapist, patient and clinic overlaps — even across the two clinics."),
     ("Packages & family sharing.", "Prepaid sessions, expiry rules, and linked spouse/child usage."),
     ("Products & follow-up.", "Herbal/supplement catalogue with price-at-sale history and follow-up list."),
     ("Built for growth.", "Every change runs through a store action — a clean seam for a future backend.")],
    note="Tech today: React + TypeScript + Capacitor (iOS/Android), client-only mock data ready for a backend.")

# --- patient experience ---
section_slide(
    "Patient experience",
    "From welcome to confirmed booking in a few taps — large readable type, calm clinical "
    "styling, and clear separation of Clinic A and Clinic B.",
    "concept-family.png", "Walkthrough · Patient")
feature_slide(
    "Patient · Onboarding", "Welcome & sign-in",
    [("welcome.png", "Welcome — Clinic A & B"), ("patient-home.png", "Patient home dashboard")],
    [("Register & verify", "Email/mobile registration with OTP/email verification before booking."),
     ("Role-based entry", "Verified patients land on a personalized home; admins are kept separate."),
     ("At-a-glance home", "Next appointment, clinic reminder, package balance and a one-tap Book shortcut.")])
feature_slide(
    "Patient · Booking", "Book in a few taps",
    [("patient-book.png", "Service → clinic → date → time")],
    [("Service first", "Choosing a service sets the duration (e.g. Physio 3h, Grounding 2h)."),
     ("Only valid slots", "Start times without room for the full duration are hidden or blocked."),
     ("Right clinic, right time", "Clear Clinic A / B choice, then date and available time."),
     ("Review before confirm", "Service, clinic, date and start–end time shown before confirming."),
     ("Double-booking blocked", "Can't take a slot just taken, or overlap an existing appointment.")])
feature_slide(
    "Patient · Appointments", "Manage visits",
    [("patient-appointments.png", "Upcoming / completed / cancelled")],
    [("Full history", "Tabs for upcoming, completed, cancelled and rescheduled visits."),
     ("Reschedule", "Move to another available slot within policy — no overlaps."),
     ("Cancel with a reason", "Pick an admin-managed reason; 'Other' allows a short note."),
     ("Cutoff respected", "Cancel/reschedule blocked past the policy cutoff (e.g. 24h)."),
     ("Balance protected", "Cancelling/rescheduling never deducts a package session.")])
feature_slide(
    "Patient · Packages & Family", "Packages & family sharing",
    [("patient-packages.png", "Package balance & expiry"), ("patient-family.png", "Link adult · add child")],
    [("See what's left", "Sessions remaining, expiry date and usage history per package."),
     ("Usage rules", "Blocked at zero balance or after expiry, even with sessions left."),
     ("Link an adult", "Spouse links by registered email/mobile — with their consent."),
     ("Add a child", "Children managed under the parent, no separate login."),
     ("Transparent usage", "See which family member used each session.")])
feature_slide(
    "Patient · Clinics & Profile", "Clinics & profile",
    [("patient-clinics.png", "Clinic A & B information"), ("patient-profile.png", "Profile & details")],
    [("Know where to go", "Clinic A and Clinic B details reduce wrong-location visits."),
     ("Keep details current", "View and update profile; required fields are validated."),
     ("Secure account", "Login, password reset after identity verification, role separation.")])

# --- admin experience ---
section_slide(
    "Admin experience",
    "Kuya Bong runs the whole business in one place — availability, bookings, packages, "
    "products and the master data that powers booking logic.",
    "concept-clinic.png", "Walkthrough · Admin")
feature_slide(
    "Admin · Command centre", "Dashboard & calendar",
    [("admin-dashboard.png", "Today, requests, follow-ups"), ("admin-calendar.png", "Availability calendar")],
    [("Day at a glance", "Today's appointments, pending requests and follow-up reminders."),
     ("Publish availability", "Open/edit/remove slots by clinic, date and time."),
     ("Duration-aware", "A 3-hour service reserves the therapist & clinic for the full window."),
     ("Safe edits", "Warned before removing a booked slot or creating overlaps.")])
feature_slide(
    "Admin · Bookings", "Appointments & manual booking",
    [("admin-appointments.png", "Filter & manage"), ("admin-manual-booking.png", "Phone/DM booking")],
    [("One inbox for all", "App, phone and DM bookings live in a single system."),
     ("Filter fast", "By clinic, date, status or therapist."),
     ("Complete a session", "Marking complete deducts one package session — only then."),
     ("Cancel on behalf", "Admin cancels with a reason + note; recorded as admin-initiated.")])
feature_slide(
    "Admin · Patients & Packages", "Patients & packages",
    [("admin-patients.png", "Search & open profiles"), ("admin-packages.png", "Create & assign")],
    [("Patient 360", "Profile, family links and appointment history in one view."),
     ("Create packages", "Set sessions and validity period for prepaid bundles."),
     ("Assign & track", "Assign to a patient; balance, expiry and usage tracked automatically."),
     ("Clear no-results", "Helpful message when a patient search finds nothing.")])
feature_slide(
    "Admin · Products", "Products & follow-up",
    [("admin-products.png", "Catalogue & prices"), ("admin-followups.png", "Follow-up list")],
    [("Manage catalogue", "Create, edit, activate/deactivate and price products."),
     ("Price snapshot", "Purchases store price-at-sale — later edits don't change history."),
     ("Record purchases", "Log a patient's purchase with date and price."),
     ("Timely follow-up", "See last-purchase dates and a follow-up list to re-engage patients.")])
feature_slide(
    "Admin · Master data", "Services, therapists & reasons",
    [("admin-services.png", "Service types & duration"), ("admin-therapists.png", "Therapists"),
     ("admin-cancellation-reasons.png", "Cancellation reasons")],
    [("Service types", "Name + standard duration drive slot logic; activate/deactivate."),
     ("Therapists", "Add Kuya Bong + occasional therapists; assign and filter by therapist."),
     ("Cancellation reasons", "Managed list (incl. 'Other' + note) for consistent records.")])

# --- requirements coverage ---
table_slide(
    "Requirements coverage", "Traceability",
    ["Blueprint / user-story area", "Where it lives in the app", "Status"],
    [["1–2  Registration, verify, login, reset", "Welcome · Register · Verify · Login · Forgot", "Built"],
     ["3  Profile & clinic info", "Patient Profile · Clinics", "Built"],
     ["4  Service-driven booking", "Patient Book (service→clinic→date→time→review)", "Built"],
     ["5  Reschedule & cancel + reasons", "Appointments · Details · Cancellation reasons", "Built"],
     ["6–7b  Availability, manual booking, admin cancel", "Calendar · Appointments · Manual booking", "Built"],
     ["8–9  Completion, packages, balance, expiry", "Appointments · Packages (patient & admin)", "Built"],
     ["10  Family linking & shared packages", "Family · Package usage", "Built"],
     ["11  Herbal/supplement catalogue & follow-up", "Products · Follow-ups", "Built"],
     ["14–18  Services, therapists, conflicts, reasons", "Services · Therapists · Cancellation reasons", "Built"]],
    note="Every mutation runs through a store action — the integration seam for the future backend.",
    col_w=[5, 5, 1.4])

bullets_slide(
    "Business rules enforced", "Rules",
    [("Verify before booking.", "Only verified accounts can reserve a session."),
     ("No double-booking.", "Therapist, patient and clinic overlaps blocked across both clinics."),
     ("Deduct on completion only.", "Cancel/reschedule never reduces a package balance."),
     ("Package guards.", "Usage blocked at zero balance or after expiry."),
     ("Price snapshot.", "Updating a catalogue price never rewrites past purchase records."),
     ("Cancel cutoff & reasons.", "Patient cutoff (e.g. 24h) and a required reason on every cancellation.")],
    note="Defaults (60-min vs service duration, auto-confirm, OTP/email) follow the blueprint and are easy to change.")

# --- COSTS ---
section_slide(
    "Operating costs",
    "What it costs to publish and run the app: one-time store fees plus recurring "
    "infrastructure. Figures are typical 2026 estimates in USD — confirm with chosen vendors.",
    "concept-cost.png", "Budget")

table_slide(
    "Developer programs & store fees", "Cost · Stores",
    ["Item", "Type", "Typical cost", "Notes"],
    [["Apple Developer Program", "Annual", "$99 / yr", "Required to publish on the App Store."],
     ["Google Play Developer", "One-time", "$25 once", "Single lifetime registration fee."],
     ["App Store review", "Per release", "$0", "Included; budget time, not money."],
     ["Code-signing / certs", "Included", "$0", "Managed inside the developer accounts."]],
    note="Year-1 store cost ≈ $124, then ≈ $99 / yr ongoing (Apple).",
    col_w=[3.4, 1.6, 2.0, 4.2])

table_slide(
    "Apple Pay & payments", "Cost · Payments",
    ["Item", "Type", "Typical cost", "Notes"],
    [["Apple Pay integration", "Per use", "$0 from Apple", "Apple charges no extra fee for Apple Pay itself."],
     ["Payment processor", "Per transaction", "≈ 2.9% + $0.30", "Stripe/Adyen etc. on card / Apple Pay / Google Pay."],
     ["In-app purchase (digital)", "Per transaction", "15–30%", "Only if selling digital goods in-app; physical services are exempt."],
     ["Merchant / gateway setup", "One-time", "$0–$50", "Often free with modern processors."]],
    note="Payments are OUT of the MVP scope (Blueprint 6.2). Booking real-world services can use external/processor payment to avoid the 15–30% IAP cut.",
    col_w=[3.2, 1.7, 2.1, 5.0])

table_slide(
    "Servers, database & domain", "Cost · Infrastructure",
    ["Item", "Starter (low traffic)", "Growth", "Notes"],
    [["Frontend / web hosting", "$0 (Netlify/Vercel free)", "$0–$20 / mo", "Static build; current demo already on Netlify."],
     ["Backend / API (BaaS)", "$0 free tier", "$25–$35 / mo", "Supabase/Firebase; scales with usage."],
     ["Database (managed Postgres)", "$0 free tier", "$15–$50 / mo", "Often bundled with the BaaS above."],
     ["File / image storage", "$0–$1 / mo", "$1–$10 / mo", "Profile photos, product images."],
     ["Domain name", "$12–$20 / yr", "$12–$20 / yr", ".app/.com; the only must-buy day one."],
     ["SMS OTP (optional)", "pay-as-you-go", "≈ $0.05–$0.08 / SMS", "Twilio; email OTP is ~free via SES/Resend."],
     ["Push notifications", "$0", "$0", "Free via Apple APNs & Firebase FCM."]],
    note="Lean launch can run at ~$0/mo infra on free tiers — only the domain (~$15/yr) is unavoidable early on.",
    col_w=[3.2, 2.8, 2.6, 4.2])

table_slide(
    "Year-1 cost summary", "Budget · Summary",
    ["Scenario", "One-time", "Recurring", "≈ Year-1 total"],
    [["Lean launch (free tiers + domain)", "$25 (Play)", "$99/yr Apple + $15/yr domain", "≈ $140"],
     ["Typical small clinic (paid BaaS)", "$25 (Play)", "~$30/mo infra + $99/yr Apple + $15/yr domain", "≈ $475"],
     ["Growth (more traffic + SMS OTP)", "$25 (Play)", "~$60/mo infra + SMS + $99 + $15", "≈ $850+"]],
    note="Excludes development effort and payment-processor % (charged only when you actually take payments).",
    col_w=[4.2, 1.8, 4.5, 2.2], highlight_last=False)

bullets_slide(
    "Cost takeaways", "Budget · Notes",
    [("Cheap to start.", "~$140 covers year one on free infra tiers + both store accounts + a domain."),
     ("Apple Pay is free to add.", "You only pay the payment processor's ~2.9% when money actually moves."),
     ("Avoid the IAP cut.", "Booking physical clinic services can settle outside in-app purchase (no 15–30%)."),
     ("Infra scales with patients.", "Stay on free tiers until real usage, then ~$30/mo handles a busy single clinic."),
     ("Recurring floor is small.", "$99/yr Apple + ~$15/yr domain are the only guaranteed ongoing costs."),
     ("Plan for the backend.", "Biggest future cost is engineering time, not hosting.")],
    note="All figures are 2026 list-price estimates in USD — validate against your chosen vendors and region.")

closing_slide()

prs.save(OUT)
print("Saved", OUT, "with", len(prs.slides._sldIdLst), "slides")
